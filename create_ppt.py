from pptx import Presentation
from pptx.util import Inches, Pt
import sys

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AI-Powered Cold Storage & Indoor Air Quality Monitoring System"
subtitle.text = "Real-time IoT Telemetry & Predictive Analytics\n\nPresented by:\nNishant Raushan Sharma\nVanisha Raj Tanwar"

# Slide 2: Intro
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Introduction & Background"
tf = slide.placeholders[1].text_frame
tf.text = "Cold storages require strict climatic control to prevent spoilage."
p = tf.add_paragraph(); p.text = "Manual monitoring is error-prone and reactive."
p = tf.add_paragraph(); p.text = "Toxic gas leaks pose immense safety and financial risks."
p = tf.add_paragraph(); p.text = "Our Solution: Real-time telemetry via IoT Edge devices and Cloud ML."

# Slide 3: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "The Core Problem"
tf = slide.placeholders[1].text_frame
tf.text = "Lack of Predictive Insights: Alarms only trigger after threshold is crossed."
p = tf.add_paragraph(); p.text = "Data Silos: Hardware alerts fail to notify stakeholders rapidly."
p = tf.add_paragraph(); p.text = "Hardware Limitations: Low-end microcontrollers struggle to run ML."

# Slide 4: Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Objectives"
tf = slide.placeholders[1].text_frame
tf.text = "Real-time Synchronization: Sub-second streaming from Edge to Web."
p = tf.add_paragraph(); p.text = "Predictive Hazard Modeling: Machine Learning to classify state."
p = tf.add_paragraph(); p.text = "Scalable Architecture: Python FastAPI & Docker for deployment."
p = tf.add_paragraph(); p.text = "Centralized Dashboard: Live visualizations of sensor data."

# Slide 5: System Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "High-Level System Overview"
tf = slide.placeholders[1].text_frame
tf.text = "Sensing Edge: ESP32 captures localized environmental data."
p = tf.add_paragraph(); p.text = "Communication Layer: HTTP POST payload over Wi-Fi."
p = tf.add_paragraph(); p.text = "Intelligent Cloud: FastAPI handles ML inference & DB storage."
p = tf.add_paragraph(); p.text = "Presentation Layer: Live dashboard rendering (WebSockets)."

# Slide 6: Input
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Input Module & Sensors"
tf = slide.placeholders[1].text_frame
tf.text = "DHT11 Sensor: Ambient Temperature (°C) & Humidity (%)."
p = tf.add_paragraph(); p.text = "MQ-2 Sensor: Detects Smoke and combustible gases."
p = tf.add_paragraph(); p.text = "MQ-135 Sensor: Air Quality (Ammonia, NOX, Benzene, CO2)."
p = tf.add_paragraph(); p.text = "Input Flow: Multi-sensor array -> ESP32 API -> JSON Payload."

# Slide 7: Output
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Output & Response Module"
tf = slide.placeholders[1].text_frame
tf.text = "Hardware Actuators: Relays, LED, and Buzzer alerts."
p = tf.add_paragraph(); p.text = "Live Visualizations: Real-time multi-dimensional charts."
p = tf.add_paragraph(); p.text = "Predictive States: Binary classification (SAFE vs DANGER)."
p = tf.add_paragraph(); p.text = "Confidence Metrics: ML probability thresholds displayed dynamically."

# Slide 8: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "System Architecture"
tf = slide.placeholders[1].text_frame
tf.text = "Tier 1 (Perception Layer): ESP32 hardware scanning environment."
p = tf.add_paragraph(); p.text = "Tier 2 (API & ML): FastAPI Web Server + Scikit-Learn Model."
p = tf.add_paragraph(); p.text = "Tier 3 (Client): JavaScript dashboard via WebSockets."
p = tf.add_paragraph(); p.text = "DevOps: SQLite logging, wrapped in Docker, deployed to Railway."

# Slide 9: Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Methodology & ML Process"
tf = slide.placeholders[1].text_frame
tf.text = "Data Synthesis: 17,500 rows of environmental sensor permutations."
p = tf.add_paragraph(); p.text = "Preprocessing: Parameters sized via robust scaling methodology."
p = tf.add_paragraph(); p.text = "Algorithm: Random Forest (Accuracy) / Logistic Regression (Speed)."
p = tf.add_paragraph(); p.text = "Inference Pipeline: Raw -> Transformation -> ML Weight Matrix -> Probability."

# Slide 10: Implementation
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Implementation & Tech Stack"
tf = slide.placeholders[1].text_frame
tf.text = "Microcontroller: C++ via Arduino IDE, ArduinoJson."
p = tf.add_paragraph(); p.text = "Backend Framework: Python 3, FastAPI, Uvicorn."
p = tf.add_paragraph(); p.text = "ML Stack: Numpy, Pandas, Scikit-Learn."
p = tf.add_paragraph(); p.text = "Frontend / Deployment: HTML5, Chart.js / Docker on Railway."

# Slide 11: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Results"
tf = slide.placeholders[1].text_frame
tf.text = "Consistent low-latency from hardware detection to global UI."
p = tf.add_paragraph(); p.text = "98.42% predictive AUC accuracy achieved via Random Forest."
p = tf.add_paragraph(); p.text = "Successful real-time plotting of 4 independent sensor streams."

# Slide 12: Advantages
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Advantages"
tf = slide.placeholders[1].text_frame
tf.text = "Edge-Cloud Synergy: Bypasses ESP32 memory constraints."
p = tf.add_paragraph(); p.text = "Platform Agnostic: Web-based dashboard from any device worldwide."
p = tf.add_paragraph(); p.text = "High Extensibility: Decoupled sensor addition logic."
p = tf.add_paragraph(); p.text = "Cost Effective: Built on affordable hardware and Open Source frameworks."

# Slide 13: Limitations
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Limitations"
tf = slide.placeholders[1].text_frame
tf.text = "Network Dependency: If Wi-Fi goes down, cloud predictions halt."
p = tf.add_paragraph(); p.text = "SQLite Constraints: Local DB lacks massive spatial scaling."
p = tf.add_paragraph(); p.text = "Sensor Calibration: MQ sensors require periodic precision tuning."

# Slide 14: Future Scope
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Future Scope"
tf = slide.placeholders[1].text_frame
tf.text = "TinyML (Edge AI): Native TensorFlow Lite Micro on ESP32."
p = tf.add_paragraph(); p.text = "Enterprise Migration: Transition to managed PostgreSQL clusters."
p = tf.add_paragraph(); p.text = "Computer Vision: Integrate ESP32-CAM for visual degradation analysis."
p = tf.add_paragraph(); p.text = "Alert Integrations: Automated SMS/WhatsApp API warnings."

# Slide 15: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
tf = slide.placeholders[1].text_frame
tf.text = "Successfully bridged physical hardware with modern cloud AI."
p = tf.add_paragraph(); p.text = "Demonstrated democratization of predictive industrial telemetry."
p = tf.add_paragraph(); p.text = "Produced an end-to-end blueprint for smart cold storage oversight."

# Slide 16: References
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "References"
tf = slide.placeholders[1].text_frame
tf.text = "Espressif Systems: ESP32 Technical Reference Manual."
p = tf.add_paragraph(); p.text = "Pedregosa et al.: Scikit-learn Machine Learning in Python."
p = tf.add_paragraph(); p.text = "FastAPI & Uvicorn Framework Documentations."
p = tf.add_paragraph(); p.text = "Chart.js Library API references."

output_path = "Cold_Storage_AI_Presentation.pptx"
prs.save(output_path)
print(f"Presentation generated successfully at {output_path}")
